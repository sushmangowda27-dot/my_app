from fastapi import FastAPI, BackgroundTasks
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pathlib import Path
import uuid
import os
import zipfile
import aiohttp
import re
from pptx import Presentation

app = FastAPI(title="Smart Project Analyzer FINAL")

# ============================
# JOB STORAGE
# ============================
jobs = {}

# ============================
# REQUEST MODEL
# ============================
class GitHubRequest(BaseModel):
    github_url: str


# ============================
# HOME
# ============================
@app.get("/")
def home():
    return {"message": "Smart Project Analyzer Running 🚀"}


# ============================
# CLEAN TEXT
# ============================
def clean_text(text: str):
    text = re.sub(r"\[.*?\]\(.*?\)", "", text)
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)
    text = re.sub(r"[*#`]", "", text)
    return text.strip()


# ============================
# GITHUB → ZIP URL
# ============================
def convert_to_zip_url(url: str):
    if url.endswith(".git"):
        url = url[:-4]

    if "github.com" not in url:
        raise ValueError("Only GitHub URLs supported")

    return url.replace("github.com", "codeload.github.com") + "/zip/main"


# ============================
# DOWNLOAD ZIP
# ============================
async def download_zip(url, path):
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as resp:
            if resp.status != 200:
                raise Exception("Failed to download repository")

            with open(path, "wb") as f:
                while True:
                    chunk = await resp.content.read(1024)
                    if not chunk:
                        break
                    f.write(chunk)


# ============================
# SAFE EXTRACT
# ============================
def safe_extract(zip_path, extract_to):
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        for member in zip_ref.namelist():
            member_path = os.path.join(extract_to, member)

            if not os.path.abspath(member_path).startswith(os.path.abspath(extract_to)):
                raise Exception("Unsafe ZIP detected")

        zip_ref.extractall(extract_to)


# ============================
# PPT GENERATOR
# ============================
def create_ppt(data, output_path):

    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("project_title", "Project")
    slide.placeholders[1].text = "Auto Generated Project Report"

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Tech Stack"
    slide.placeholders[1].text = "\n".join(data.get("tech_stack", []))

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    slide.placeholders[1].text = data.get("readme_preview", "")[:1000]

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Stats"
    slide.placeholders[1].text = (
        f"Total Files: {data.get('total_files')}\n"
        f"Tech Stack: {', '.join(data.get('tech_stack', []))}"
    )

    prs.save(output_path)


# ============================
# BACKGROUND PROCESS
# ============================
def process_repo(job_id, github_url):

    try:
        jobs[job_id]["status"] = "downloading"

        base_path = Path("backend/storage/repos") / job_id
        base_path.mkdir(parents=True, exist_ok=True)

        zip_path = base_path / "repo.zip"

        import asyncio
        zip_url = convert_to_zip_url(github_url)
        asyncio.run(download_zip(zip_url, str(zip_path)))

        jobs[job_id]["status"] = "extracting"
        safe_extract(str(zip_path), str(base_path))
        os.remove(zip_path)

        extracted_folder = next(base_path.iterdir())

        jobs[job_id]["status"] = "analyzing"

        files = []
        for f in extracted_folder.rglob("*"):
            if ".git" not in str(f) and f.is_file():
                files.append(f.name)

        readme = ""
        for f in extracted_folder.rglob("README*"):
            readme = f.read_text(errors="ignore")
            break

        lines = [l.strip() for l in readme.splitlines() if l.strip()]

        project_title = "Unknown Project"
        for line in lines:
            if line.startswith("#"):
                project_title = clean_text(line.replace("#", ""))
                break

        if project_title == "Unknown Project" and lines:
            project_title = clean_text(lines[0])

        tech_stack = []
        if "package.json" in files:
            tech_stack.append("Node.js")
        if "requirements.txt" in files:
            tech_stack.append("Python")
        if "pubspec.yaml" in files:
            tech_stack.append("Flutter")
        if "pom.xml" in files:
            tech_stack.append("Java")

        result = {
            "project_title": project_title,
            "total_files": len(files),
            "tech_stack": tech_stack,
            "sample_files": files[:20],
            "readme_preview": readme[:3000]
        }

        jobs[job_id]["status"] = "completed"
        jobs[job_id]["result"] = result

    except Exception as e:
        jobs[job_id]["status"] = "failed"
        jobs[job_id]["error"] = str(e)


# ============================
# SUBMIT GITHUB
# ============================
@app.post("/submit-github")
def submit_github(data: GitHubRequest, background_tasks: BackgroundTasks):

    job_id = str(uuid.uuid4())

    jobs[job_id] = {
        "status": "queued",
        "result": None,
        "error": None
    }

    background_tasks.add_task(process_repo, job_id, data.github_url)

    return {"job_id": job_id, "status": "started"}


# ============================
# STATUS API (FIXED)
# ============================
@app.get("/status/{job_id}")
def status(job_id: str):

    job = jobs.get(job_id)

    if not job:
        return {"status": "error", "message": "Invalid job id"}

    return {
        "job_id": job_id,
        "status": job["status"],
        "error": job.get("error")
    }


# ============================
# RESULT API (FIXED)
# ============================
@app.get("/result/{job_id}")
def result(job_id: str):

    job = jobs.get(job_id)

    if not job:
        return {"status": "error", "message": "Invalid job id"}

    if job["status"] != "completed":
        return {
            "status": job["status"],
            "message": "Result not ready yet"
        }

    return job["result"]


# ============================
# PPT GENERATE API (SAFE)
# ============================
@app.get("/generate-ppt/{job_id}")
def generate_ppt(job_id: str):

    job = jobs.get(job_id)

    if not job:
        return {"status": "error", "message": "Invalid job id"}

    if job["status"] != "completed":
        return {"status": "not_ready", "message": "Analysis not completed yet"}

    ppt_dir = Path("backend/storage/ppt")
    ppt_dir.mkdir(parents=True, exist_ok=True)

    ppt_path = ppt_dir / f"{job_id}.pptx"

    create_ppt(job["result"], ppt_path)

    return {
        "status": "success",
        "ppt_file": str(ppt_path)
    }


# ============================
# DOWNLOAD PPT (FINAL FIX)
# ============================
@app.get("/download-ppt/{job_id}")
def download_ppt(job_id: str):

    ppt_path = Path(f"backend/storage/ppt/{job_id}.pptx")

    if not ppt_path.exists():
        return {
            "status": "error",
            "message": "PPT not found. Generate it first."
        }

    return FileResponse(
        path=ppt_path,
        filename=f"{job_id}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )